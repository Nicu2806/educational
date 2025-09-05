import streamlit as st
from g4f.client import Client
from g4f.errors import RateLimitError, ProviderNotFoundError, ModelNotFoundError
import pandas as pd
import re
import json
import os
from docx import Document as DocxDocument
from io import BytesIO
import pdfplumber
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import inspect

# --- Initial Data (can be overwritten) ---
DEFAULT_COMPETENCIES_SPECIFIC = {
    "C1": "Knowledge of fundamental concepts of computer graphics (image types, standards, graphic systems and libraries).",
    "C2": "Application of 2D and 3D geometric transformations (translation, scaling, rotation, mirroring, shearing, composite transformations).",
    "C3": "Use of homogeneous coordinates and transformation matrices for manipulating graphic scenes.",
    "C4": "Implementation of fundamental drawing algorithms (DDA, Bresenham for segments, circles, ellipses).",
    "C5": "Modeling and visualizing 3D scenes using projections (parallel and perspective).",
    "C6": "Application of realistic image rendering techniques (shading, lighting, transparency, reflection)."
}

MATERIAL_CATEGORIES = [
    "Curriculum", "Textbooks", "Support Materials (pdf, ppt)",
    "Worksheets", "Evaluation Sheets", "Interdisciplinary Projects"
]

COVERAGE_LEGEND = {
    "Complete Coverage ✅": "✅",
    "Partial Coverage 🤔": "🤔",
    "Missing/Not Covered ❌": "❌"
}
COVERAGE_LEGEND_EN_KEYS_FOR_AI = {
    "Complete": "✅",
    "Partial": "🤔",
    "Missing": "❌"
}

# --- Data Management Helper Functions ---

def get_curriculum_context():
    """Loads the aggregated text content from the user's session state."""
    if 'processed_data' not in st.session_state or not st.session_state.processed_data:
        return None, 0

    data = st.session_state.processed_data
    
    full_text = "\n\n--- FILE SEPARATOR ---\n\n".join(
        f"Content from file '{path}':\n{content}"
        for path, content in data.items()
    )
    return full_text, len(data)

def extract_text_from_file(uploaded_file):
    """
    Extracts text from a single uploaded file object by checking its filename extension.
    This is more reliable than using the browser-provided MIME type.
    """
    if uploaded_file is None:
        return ""
    
    file_name = uploaded_file.name
    ext = os.path.splitext(file_name)[1].lower()
    text = ""
    
    try:
        uploaded_file.seek(0)
        
        if ext == ".pdf":
            with pdfplumber.open(uploaded_file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        elif ext == ".docx":
            doc = DocxDocument(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == ".pptx":
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in [".txt", ".md"]:
            text = uploaded_file.getvalue().decode("utf-8")
        else:
            st.warning(f"File extension '{ext}' for '{file_name}' is not supported for direct extraction.")
            
    except Exception as e:
        st.error(f"An error occurred while reading the file '{file_name}': {e}")

    return text

# --- AI Helper Functions ---
@st.cache_resource
def init_ai_service_client():
    """Initializes the AI client."""
    try:
        client = Client()
        return client
    except Exception as e:
        st.error(f"Critical error initializing AI client: {e}")
        return None

@st.cache_resource
def get_functional_models():
    """
    Tests all available models to find which ones are functional.
    This function is cached to run only once at the start.
    """
    functional_models = []
    all_model_classes = inspect.getmembers(g4f.models, inspect.isclass)
    test_prompt = "Scrie o singură propoziție despre un robot."
    
    # Prioritize the user's requested models for testing
    priority_models = ["DeepInfra", "LambdaChat", "OIVSCodeSer0501", "WeWordle", "Yqcloud"]
    
    # Create a set of all model names to avoid re-testing
    all_model_names = {name for name, _ in all_model_classes}
    
    # Test priority models first
    for model_name in priority_models:
        if model_name in all_model_names:
            try:
                print(f"Testing priority model: {model_name}")
                response = g4f.ChatCompletion.create(
                    model=model_name,
                    messages=[{"role": "user", "content": test_prompt}],
                    stream=False,
                    timeout=10
                )
                if response:
                    functional_models.append(model_name)
            except Exception as e:
                print(f"Model {model_name} failed: {e}")
                
    # Test the rest of the models
    for name, model_class in all_model_classes:
        if name not in priority_models and not name.startswith('_'):
            try:
                print(f"Testing other model: {name}")
                response = g4f.ChatCompletion.create(
                    model=name,
                    messages=[{"role": "user", "content": test_prompt}],
                    stream=False,
                    timeout=10
                )
                if response:
                    functional_models.append(name)
            except Exception as e:
                print(f"Model {name} failed: {e}")
    
    if not functional_models:
        st.error("No functional AI models were found. Please try again later.")
        
    return functional_models

def process_direct_with_ai_service(user_input_text, system_prompt, ai_client_instance, generation_params=None):
    """
    Sends a request to the AI service and returns the response.
    It now automatically tries a list of models until it finds a working one.
    """
    params = generation_params.copy() if generation_params else {}
    
    if "messages_override" in params:
        messages_to_send = params.pop("messages_override")
    else:
        messages_to_send = []
        if system_prompt:
            messages_to_send.append({"role": "system", "content": system_prompt})
        if user_input_text:
            messages_to_send.append({"role": "user", "content": user_input_text})

    if not messages_to_send:
        return "Internal Error: No input provided for AI."
    
    if ai_client_instance is None:
        return "AI Service is unavailable."

    available_models_to_try = get_functional_models()
    
    used_model = None
    response = None

    for model_name in available_models_to_try:
        params['model'] = model_name
        try:
            print(f"Încercăm modelul: {model_name}")
            response = ai_client_instance.chat.completions.create(
                messages=messages_to_send,
                stream=False,
                timeout=180,
                **params
            )
            if response.choices and response.choices[0].message and response.choices[0].message.content:
                used_model = model_name
                break # Ieșim din buclă dacă am găsit un răspuns valid
        except Exception as e:
            print(f"Eroare cu modelul {model_name}: {e}")
            continue

    if not used_model:
        return "Nu a fost găsit niciun model funcțional."

    if response.choices and response.choices[0].message and response.choices[0].message.content:
        return response.choices[0].message.content.strip()
    
    st.warning(f"Unexpected AI response structure: {response}")
    return "The AI service did not return valid content."

# --- Formatting and Document Generation Helper Functions ---

def format_cell_for_custom_display(cell_text_content):
    """Formats cell content with icons and colors for Streamlit display."""
    if pd.isna(cell_text_content) or not str(cell_text_content).strip() or "Missing" in str(cell_text_content):
        return f"<span style='color:red;'>{COVERAGE_LEGEND_EN_KEYS_FOR_AI['Missing']}</span>&nbsp;<em>Missing</em>"

    cell_text_normalized = str(cell_text_content).replace("\n", " ")
    
    items_html = []
    extracted_pairs = re.findall(r"([^,(]+?(?:\.[\w\d]+)?)\s*\(.*?(✅|🤔|❌).*?\)", cell_text_normalized)

    if not extracted_pairs:
        return f"<em>{cell_text_content}</em>"

    for item_name, status_symbol in extracted_pairs:
        item_name = item_name.strip()
        item_symbol = status_symbol.strip()
        item_color = "grey"
        if item_symbol == COVERAGE_LEGEND_EN_KEYS_FOR_AI["Complete"]: item_color = "green"
        elif item_symbol == COVERAGE_LEGEND_EN_KEYS_FOR_AI["Partial"]: item_color = "orange"
        
        items_html.append(f"<div style='margin-bottom: 5px;'><span style='color:{item_color}; font-size: 1.1em;'>{item_symbol}</span>&nbsp;{item_name}</div>")
        
    return "".join(items_html)

def create_document_word(content, title="Generated Document"):
    """Creates a Word document from the provided content with basic styling."""
    document = DocxDocument()
    document.add_heading(title, 0)
    for line in content.split('\n'):
        line = line.strip()
        if line.startswith('### '):
            document.add_heading(line.replace('### ', ''), level=3)
        elif line.startswith('## '):
            document.add_heading(line.replace('## ', ''), level=2)
        elif line.startswith('# '):
            document.add_heading(line.replace('## ', ''), level=1)
        elif re.match(r'^\d+\.', line):
            document.add_paragraph(line, style='List Number')
        elif line.startswith('* ') or line.startswith('- '):
            document.add_paragraph(line, style='List Bullet')
        elif line.startswith('**') and line.endswith('**'):
            p = document.add_paragraph()
            p.add_run(line.strip('**')).bold = True
        else:
            document.add_paragraph(line)
            
    buffer = BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer

def create_presentation_from_text(slide_text):
    """Creates a PowerPoint presentation from Markdown-like formatted text."""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    title_font = 'Calibri'
    content_font = 'Calibri'
    title_font_size = Pt(44)
    content_font_size = Pt(22)
    title_color = RGBColor(0x0A, 0x42, 0x6E)
    content_color = RGBColor(0x21, 0x21, 0x21)

    # Split text into slide blocks based on "## Slide"
    slide_blocks = re.split(r'\n##\s+Slide.*', slide_text)

    for block in slide_blocks:
        if not block.strip():
            continue

        lines = [line.strip() for line in block.strip().split('\n') if line.strip()]
        if not lines:
            continue
        
        # The first line is the title
        title_text = lines[0]
        content_lines = lines[1:]

        # Use a title and content layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Set slide title
        title_shape = slide.shapes.title
        title_shape.text = title_text
        p_title = title_shape.text_frame.paragraphs[0]
        p_title.font.name = title_font
        p_title.font.size = title_font_size
        p_title.font.bold = True
        p_title.font.color.rgb = title_color

        # Set slide content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # Clear default text

        for line in content_lines:
            if not line:
                continue
            
            p = tf.add_paragraph()
            
            # Handle bullet points (level 1 and 2)
            if line.startswith('- ') or line.startswith('* '):
                p.text = line[2:]
                p.level = 1
            elif re.match(r'^\s*-\s|\s*\*\s', line): # indented bullet
                p.text = re.sub(r'^\s*[-\*]\s*', '', line)
                p.level = 2
            else:
                p.text = line
                p.level = 0
            
            p.font.name = content_font
            p.font.size = content_font_size
            p.font.color.rgb = content_color
            p.space_after = Pt(12)
            
        # Remove the initial empty paragraph if it exists
        if tf.paragraphs and not tf.paragraphs[0].text and len(tf.paragraphs) > 1:
            tf._element.remove(tf.paragraphs[0]._element)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
