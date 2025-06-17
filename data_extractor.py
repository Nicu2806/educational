# praga/data_extractor.py

import os
import json
from docx import Document
import pdfplumber
from pptx import Presentation
import streamlit as st

def extract_text_from_docx(path):
    """Extracts text from a .docx file."""
    try:
        doc = Document(path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        st.warning(f"Could not read DOCX file {os.path.basename(path)}: {e}")
        return ""

def extract_text_from_pdf(path):
    """Extracts text from a .pdf file."""
    text = ""
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text
    except Exception as e:
        st.warning(f"Could not read PDF file {os.path.basename(path)}: {e}")
        return ""

def extract_text_from_pptx(path):
    """Extracts text from a .pptx file."""
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.warning(f"Could not read PPTX file {os.path.basename(path)}: {e}")
        return ""

def extract_from_folder(root_folder):
    """
    Recursively extracts text from all supported files in a folder
    and returns a dictionary with file paths and their content.
    """
    data = {}
    supported_extensions = {".docx", ".pdf", ".pptx", ".txt", ".md"}
    
    for dirpath, _, files in os.walk(root_folder):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext not in supported_extensions:
                continue

            full_path = os.path.join(dirpath, file)
            text = ""
            
            if ext == ".docx":
                text = extract_text_from_docx(full_path)
            elif ext == ".pdf":
                text = extract_text_from_pdf(full_path)
            elif ext == ".pptx":
                text = extract_text_from_pptx(full_path)
            elif ext in [".txt", ".md"]:
                try:
                    with open(full_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                except Exception as e:
                    st.warning(f"Could not read text file {file}: {e}")
                    text = ""
            
            if text.strip():
                # Use relative path for cleaner keys
                relative_path = os.path.relpath(full_path, root_folder)
                data[relative_path] = text
                
    return data

def process_folder_and_save_json(folder_path, output_dir="output"):
    """
    Processes a folder, extracts text, and saves it to a JSON file.
    Returns the path to the JSON file or None if failed.
    """
    if not os.path.isdir(folder_path):
        st.error(f"The provided path is not a valid folder: {folder_path}")
        return None
        
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    st.info(f"Starting text extraction from folder: {folder_path}...")
    extracted_data = extract_from_folder(folder_path)
    
    if not extracted_data:
        st.warning("No text could be extracted from the supported files in the folder.")
        return None
        
    json_path = os.path.join(output_dir, "data_content.json")
    
    try:
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(extracted_data, f, ensure_ascii=False, indent=4)
        st.success(f"Extracted content from {len(extracted_data)} files and saved to {json_path}")
        return json_path
    except Exception as e:
        st.error(f"Failed to save data to JSON file: {e}")
        return None